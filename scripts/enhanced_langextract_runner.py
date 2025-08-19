#!/usr/bin/env python3
"""
Enhanced LangExtract runner script with comprehensive file support.
This script bridges the Node.js backend with the Python LangExtract library.
"""

import sys
import json
import os
import tempfile
import time
from typing import List, Dict, Any, Optional
import traceback

# File processing imports
import requests
from pathlib import Path
import mimetypes

# Document processing imports
try:
    import docx
    from docx import Document
except ImportError:
    docx = None

try:
    import pdfplumber
    import PyPDF2
except ImportError:
    pdfplumber = None
    PyPDF2 = None

try:
    import openpyxl
    from openpyxl import load_workbook
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from bs4 import BeautifulSoup

def extract_text_from_file(file_path: str, mime_type: str = None) -> str:
    """Extract text from various file formats."""
    if not mime_type:
        mime_type, _ = mimetypes.guess_type(file_path)
    
    file_ext = Path(file_path).suffix.lower()
    
    try:
        # Text files
        if mime_type and 'text' in mime_type or file_ext in ['.txt', '.md', '.csv']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        
        # PDF files
        elif file_ext == '.pdf' and pdfplumber:
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text.strip()
        
        # Word documents
        elif file_ext in ['.docx', '.doc'] and docx:
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)
        
        # Excel files
        elif file_ext in ['.xlsx', '.xls'] and openpyxl:
            workbook = load_workbook(file_path, data_only=True)
            text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"Sheet: {sheet_name}")
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else '' for cell in row]
                    if any(row_text):
                        text.append('\t'.join(row_text))
            return '\n'.join(text)
        
        # PowerPoint files
        elif file_ext in ['.pptx', '.ppt'] and Presentation:
            prs = Presentation(file_path)
            text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        
        # HTML files
        elif file_ext in ['.html', '.htm'] or (mime_type and 'html' in mime_type):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                return soup.get_text(separator='\n', strip=True)
        
        # JSON files
        elif file_ext == '.json':
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        
        # Fallback to text reading
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
                
    except Exception as e:
        raise ValueError(f"Failed to extract text from {file_ext} file: {str(e)}")

def download_and_extract_from_url(url: str) -> str:
    """Download and extract text from a URL."""
    try:
        # Download the content
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()
        
        content_type = response.headers.get('content-type', '').lower()
        
        # Handle different content types
        if 'text/html' in content_type:
            soup = BeautifulSoup(response.content, 'html.parser')
            return soup.get_text(separator='\n', strip=True)
        elif 'text/' in content_type:
            return response.text
        elif 'application/pdf' in content_type and pdfplumber:
            # Save PDF temporarily and extract
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
                tmp_file.write(response.content)
                tmp_file.flush()
                return extract_text_from_file(tmp_file.name, 'application/pdf')
        else:
            # Try to extract as text
            return response.text
            
    except Exception as e:
        raise ValueError(f"Failed to download or extract text from URL: {str(e)}")

def main():
    """Main function to run LangExtract based on provided configuration."""
    start_time = time.time()
    
    try:
        # Read configuration from stdin
        config_json = sys.stdin.read()
        config = json.loads(config_json)
        
        # Validate required fields
        required_fields = ['inputText', 'promptDescription', 'examples', 'modelId']
        for field in required_fields:
            if field not in config:
                raise ValueError(f"Missing required field: {field}")
        
        # Handle file input if provided
        input_text = config['inputText']
        file_path = config.get('filePath')
        
        if file_path and os.path.exists(file_path):
            try:
                input_text = extract_text_from_file(file_path)
                if not input_text.strip():
                    raise ValueError("No text could be extracted from the file")
                
                # For text extraction only, return the text directly
                if config['promptDescription'] == "Extract text content":
                    output = {
                        'success': True,
                        'extractedText': input_text,
                        'metadata': {
                            'inputLength': len(input_text),
                            'processingTime': (time.time() - start_time) * 1000
                        }
                    }
                    print(json.dumps(output))
                    return
                    
            except Exception as e:
                raise ValueError(f"File processing failed: {str(e)}")
        
        # Handle URL input
        elif input_text.startswith(('http://', 'https://')):
            try:
                input_text = download_and_extract_from_url(input_text)
                if not input_text.strip():
                    raise ValueError("No text could be extracted from the URL")
            except Exception as e:
                raise ValueError(f"URL processing failed: {str(e)}")
        
        # Set API key from environment if not provided in config
        api_key = config.get('apiKey') or os.environ.get('GEMINI_API_KEY')
        if not api_key:
            raise ValueError("No API key provided. Set GEMINI_API_KEY environment variable.")
        
        # Set the API key in environment for LangExtract
        os.environ['GEMINI_API_KEY'] = api_key
        
        # Import LangExtract after setting environment variables
        try:
            import langextract as lx
        except ImportError as e:
            raise ValueError(f"LangExtract library not available: {str(e)}")
        
        # Prepare examples
        examples = []
        for example_data in config['examples']:
            extractions = []
            for extraction_data in example_data['extractions']:
                extraction = lx.data.Extraction(
                    extraction_class=extraction_data['extraction_class'],
                    extraction_text=extraction_data['extraction_text'],
                    attributes=extraction_data.get('attributes', {})
                )
                extractions.append(extraction)
            
            example = lx.data.ExampleData(
                text=example_data['text'],
                extractions=extractions
            )
            examples.append(example)
        
        # Run extraction with error handling
        try:
            result = lx.extract(
                text_or_documents=input_text,
                prompt_description=config['promptDescription'],
                examples=examples,
                model_id=config['modelId'],
                extraction_passes=config.get('extractionPasses', 1),
                max_workers=config.get('maxWorkers', 5),
                max_char_buffer=config.get('maxCharBuffer', 10000)
            )
        except Exception as e:
            raise ValueError(f"LangExtract processing failed: {str(e)}")
        
        # Calculate processing time
        processing_time = (time.time() - start_time) * 1000  # Convert to milliseconds
        
        # Convert result to JSON-serializable format
        output = {
            'success': True,
            'extractions': [],
            'metadata': {
                'totalExtractions': 0,
                'uniqueClasses': 0,
                'processingTime': processing_time,
                'inputLength': len(input_text),
                'averageConfidence': 0
            }
        }
        
        if hasattr(result, 'extractions') and result.extractions:
            confidences = []
            for extraction in result.extractions:
                extraction_dict = {
                    'extraction_class': extraction.extraction_class,
                    'extraction_text': extraction.extraction_text,
                    'attributes': extraction.attributes,
                    'position_start': getattr(extraction, 'position_start', None),
                    'position_end': getattr(extraction, 'position_end', None),
                    'confidence': getattr(extraction, 'confidence', None)
                }
                output['extractions'].append(extraction_dict)
                
                if extraction_dict['confidence'] is not None:
                    confidences.append(extraction_dict['confidence'])
            
            # Calculate metadata
            output['metadata']['totalExtractions'] = len(output['extractions'])
            output['metadata']['uniqueClasses'] = len(set(e['extraction_class'] for e in output['extractions']))
            if confidences:
                output['metadata']['averageConfidence'] = sum(confidences) / len(confidences)
        
        # Save results to temporary file for visualization
        if output['extractions']:
            with tempfile.NamedTemporaryFile(mode='w', suffix='.jsonl', delete=False) as f:
                json.dump({
                    'text': input_text,
                    'extractions': output['extractions']
                }, f)
                output['visualization_file'] = f.name
        
        print(json.dumps(output))
        
    except Exception as e:
        error_output = {
            'success': False,
            'error': str(e),
            'type': type(e).__name__,
            'traceback': traceback.format_exc()
        }
        print(json.dumps(error_output))
        sys.exit(1)

if __name__ == '__main__':
    main()